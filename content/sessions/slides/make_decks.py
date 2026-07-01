#!/usr/bin/env python3
"""Generate the two Radiant Women's Circle PowerPoint decks in Tarisha's brand style.
Simple, full-bleed coloured backgrounds from the locked palette, big elegant
headings, one sentence or a few bullets per slide, navy/brand footer."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---- Brand palette (locked) ----
CREAM    = RGBColor(0xF5, 0xEF, 0xE3)
IVORY    = RGBColor(0xFB, 0xF7, 0xEE)
BURGUNDY = RGBColor(0x6E, 0x1A, 0x2E)
NAVY     = RGBColor(0x1F, 0x2A, 0x44)
RUST     = RGBColor(0x9E, 0x4A, 0x2A)
COPPER   = RGBColor(0xC7, 0x5D, 0x3D)
GOLD     = RGBColor(0xC2, 0xA4, 0x6D)
NEARBLK  = RGBColor(0x15, 0x11, 0x0D)

# Fonts
HEAD_FONT = "Marcellus"   # locked brand headline font (sexualempowermentforwomen)
BODY_FONT = "Lora"        # locked brand body font

# text colour + accent for each background
SCHEME = {
    "cream":    (CREAM,    NEARBLK, COPPER),
    "ivory":    (IVORY,    NEARBLK, COPPER),
    "burgundy": (BURGUNDY, CREAM,   GOLD),
    "navy":     (NAVY,     CREAM,   GOLD),
    "rust":     (RUST,     CREAM,   GOLD),
    "copper":   (COPPER,   CREAM,   CREAM),
    "gold":     (GOLD,     NEARBLK, BURGUNDY),
}

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs


def add_slide(prs, bg, title, body=None, eyebrow=None, big=False, footer=True):
    bg_col, txt_col, accent = SCHEME[bg]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # full-bleed background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_col

    left = Inches(0.9)
    width = Inches(11.5)
    top = Inches(1.1)

    # eyebrow (small caps label)
    if eyebrow:
        eb = slide.shapes.add_textbox(left, Inches(0.6), width, Inches(0.5))
        p = eb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = eyebrow.upper()
        r.font.name = BODY_FONT; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = accent
        # letterspacing-ish: spaces already added by caller if wanted

    # title
    tb = slide.shapes.add_textbox(left, top, width, Inches(3.2))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.name = HEAD_FONT
    r.font.size = Pt(54 if big else 40)
    r.font.color.rgb = txt_col

    # accent divider under the title
    line = slide.shapes.add_shape(1, left, top + Inches(3.25 if big else 2.0),
                                  Inches(1.6), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = accent
    line.line.fill.background()
    line.shadow.inherit = False

    # body: sentence or bullets
    if body:
        by = slide.shapes.add_textbox(left, top + Inches(3.6 if big else 2.35),
                                      width, Inches(3.2))
        bf = by.text_frame; bf.word_wrap = True
        items = body if isinstance(body, list) else [body]
        for i, item in enumerate(items):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            bullet = isinstance(body, list) and len(body) > 1
            run = para.add_run()
            run.text = (u"•  " + item) if bullet else item
            run.font.name = BODY_FONT
            run.font.size = Pt(23 if bullet else 27)
            run.font.italic = not bullet
            run.font.color.rgb = txt_col
            para.space_after = Pt(12)

    # footer
    if footer:
        ft = slide.shapes.add_textbox(left, Inches(6.85), width, Inches(0.4))
        fp = ft.text_frame.paragraphs[0]
        fr = fp.add_run(); fr.text = "SEXUALEMPOWERMENTFORWOMEN.COM"
        fr.font.name = BODY_FONT; fr.font.size = Pt(11)
        fr.font.color.rgb = accent
    return slide


# =====================================================================
# SESSION 1 — Coming Home to Yourself
# =====================================================================
s1 = new_deck()

add_slide(s1, "burgundy", "Coming Home to Yourself",
          eyebrow="Radiant Women's Circle  ·  Session 1",
          body="Finding what you want, need and desire.", big=True, footer=True)

add_slide(s1, "cream", "Tonight, we find her.",
          body="Your needs. Your wants. Your deepest desires. We start by coming home to yourself.")

add_slide(s1, "navy", "Our circle",
          eyebrow="How we hold this space",
          body=["What's shared here stays here.",
                "We witness. We don't fix or rescue.",
                "You can always pass.",
                "There's no such thing as too much in this room."])

add_slide(s1, "rust", "Soften your jaw.",
          body="Drop your shoulders. Unclench your belly. She doesn't live in the tension. She lives in the softening.")

add_slide(s1, "cream", "Three words we blur together",
          body=["Needs — what you can't be well without.",
                "Wants — the specific things that meet a need.",
                "Desires — the deep, alive pull. Your compass."])

add_slide(s1, "gold", "You're not too much.",
          body="You've been gathering everyone else's needs for so long you lost the thread of your own. Tonight we pick it back up.")

add_slide(s1, "burgundy", "The invisible woman",
          eyebrow="Claire Zammit",
          body="You are not overlooked by accident.", big=True)

add_slide(s1, "navy", "How we disappear",
          eyebrow="The I'm Invisible pattern",
          body=["My attention is on everyone but me.",
                "I assume you can see what I never said.",
                "I serve until I'm empty.",
                "I wait, quietly, to be discovered."])

add_slide(s1, "copper", "The loop",
          body="Believe you're invisible → disappear yourself → they don't see you → you feel invisible. The belief builds its own proof.")

add_slide(s1, "cream", "The hard truth, said gently",
          body="You have taught the people around you exactly how much of you to see. A habit, not a fact. And tonight we start practising the opposite.")

add_slide(s1, "gold", "I see myself. I am present.",
          eyebrow="The shift  ·  presencing",
          body="I stop waiting to be seen. I turn toward myself, and I make myself visible. It is my destiny to be visible.")

add_slide(s1, "rust", "Breakout · How I make myself invisible",
          eyebrow="In pairs",
          body=["One way I make myself invisible is...",
                "One way I've trained people not to see my needs is...",
                "The story underneath might be... (invisible / not enough / too much / a burden)",
                "Partner: just witness. Then swap."])

add_slide(s1, "burgundy", "What I don't want",
          eyebrow="Part one",
          body="The no is the doorway. Your body already knows it.", big=True)

add_slide(s1, "cream", "Name the no",
          eyebrow="Journal",
          body=["What I'm tired of pretending is...",
                "What I'm done tolerating is...",
                "What I no longer want is..."])

add_slide(s1, "navy", "What I do want",
          eyebrow="Part two",
          body="Every no is holding a yes behind it.", big=True)

add_slide(s1, "cream", "Turn it into a yes",
          eyebrow="Journal, then say it out loud",
          body=["The kind of touch I'm craving is...",
                "A way I want to receive more is...",
                "If I trusted myself completely, I would..."])

add_slide(s1, "burgundy", "My deepest desire",
          eyebrow="Part three",
          body="Underneath the want is the desire your whole life is organised around.", big=True)

add_slide(s1, "rust", "Drop underneath",
          body="And if I had that, what would it give me? Keep asking. Reach for the true answer, not the impressive one.")

add_slide(s1, "gold", "Claim it",
          body="I'm a radiant woman, and I desire... The circle answers: You're allowed.")

add_slide(s1, "navy", "I take my place",
          eyebrow="The power statement",
          body="I see myself. I am present to my own feelings, needs and desires. It is my destiny to be visible, and I take my rightful place.")

add_slide(s1, "burgundy", "Claim your power",
          body="I'm powerful. I can create life. I can create the life I desire. The power is in my hands.")

add_slide(s1, "cream", "This week",
          eyebrow="Take home",
          body=["Say one clean want a day, out loud.",
                "Each night, whisper your deepest desire to yourself.",
                "You're teaching your body that your desire is safe."])

s1_path = "content/sessions/slides/Session-1-Coming-Home-to-Yourself.pptx"

# =====================================================================
# SESSION 2 — Speaking It So You're Heard
# =====================================================================
s2 = new_deck()

add_slide(s2, "navy", "Speaking It So You're Heard",
          eyebrow="Radiant Women's Circle  ·  Session 2",
          body="Asking in a way that inspires.", big=True)

add_slide(s2, "cream", "Last time we found her.",
          body="Tonight, we give her a voice. Not a complaint. Not a hint. A desire, spoken so it can be met.")

add_slide(s2, "burgundy", "Reconnect",
          body="One desire I found last time was...")

add_slide(s2, "navy", "How you're wired",
          eyebrow="Alison Armstrong",
          body=["You live in diffuse awareness — aware of everything at once.",
                "He lives more in single focus — one thing at a time.",
                "So when you share, he hears information, not a request."])

add_slide(s2, "rust", "Sharing isn't asking.",
          body="A hundred diffuse hints land as noise. One clear want lands as a request.")

add_slide(s2, "cream", "A real request",
          body=["Specific — he can picture it.",
                "A by-when — it lands in time.",
                "Room for a yes or a no — a request, not a demand."])

add_slide(s2, "copper", "Breakout · Sharing vs Asking",
          eyebrow="In pairs",
          body=["Share it the messy way first.",
                "Then say it in one clean sentence: I want... by...",
                "Partner reflects back the clear request. Swap."])

add_slide(s2, "burgundy", "The Great Ask",
          eyebrow="Four steps",
          body=["Describe exactly what you want.",
                "Say what it would give you.",
                "Say what's in it for him.",
                "Ask: what do you need from me to give me that?"])

add_slide(s2, "navy", "Frustration is a gap, not a fact.",
          body="The distance between an expectation you never spoke and the influence you didn't use. An unspoken expectation is a resentment in waiting.")

add_slide(s2, "rust", "Breakout · The Great Ask, out loud",
          eyebrow="In pairs",
          body=["Build your ask in all four steps.",
                "Partner receives it warmly, in character.",
                "No apologising. Say it like it's allowed. Swap."])

add_slide(s2, "burgundy", "Inspire, don't demand.",
          body="A demand makes him wrong and he stops trying. An invitation hands him a way to be your hero.", big=True)

add_slide(s2, "gold", "Appreciation unlocks generosity.",
          body="Notice what he already does before you ask for the next thing. Then let him see it landed. That's the fuel.")

add_slide(s2, "cream", "The loop that works",
          body=["Clear invitation → he gives.",
                "You receive it and show it landed.",
                "He wants to give more.",
                "Same man. Same marriage. You changed how you asked."])

add_slide(s2, "copper", "Receive like a queen.",
          body="Stop deflecting the gift. Let the compliment land. You can't ask for intimacy from an empty tank.")

add_slide(s2, "rust", "Breakout · Demand into invitation",
          eyebrow="In pairs",
          body=["Say the want as a demand: You never...",
                "Then as an invitation: I'd so love it if you would...",
                "Show him it would land: ...and it would make me feel...",
                "Partner: which made you want to give? Swap."])

add_slide(s2, "burgundy", "Claim it",
          body="This week, the one clean want I'm going to actually say is... The circle answers: You're allowed.")

add_slide(s2, "cream", "This week",
          eyebrow="Take home",
          body=["One clean, specific, appreciation-first request a day.",
                "Watch the generosity come back.",
                "You're teaching the people who love you how to succeed at loving you."])

s2_path = "content/sessions/slides/Session-2-Speaking-It-So-Youre-Heard.pptx"

os.makedirs("content/sessions/slides", exist_ok=True)
s1.save(s1_path)
s2.save(s2_path)
print("saved:", s1_path, len(s1.slides.__iter__.__self__._sldIdLst), "slides")
print("saved:", s2_path)
print("S1 slides:", len(s1.slides._sldIdLst))
print("S2 slides:", len(s2.slides._sldIdLst))
